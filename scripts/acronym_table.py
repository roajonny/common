# File :            acronym_table.py
# Title :           acronym_table
# 
# Author(s) :       Jonathan Roa
# 
# Description :     Generates an .ppt table from an text file 
#                   following the format (see "acronyms.txt"):
#
#                   <ACRONYM>, <DESCRIPTION,
#
# Dependencies :    python-pptx (v1.0.0)
# 
# Revisions 
# 
# Date        Name            REV#        Description 
# ----------  --------------- ----------- --------------------
# (02/21/25)  Jonathan Roa    1.0         Initial Revision

from pptx import Presentation
from pptx.util import Inches

# Have the script target the acronyms file
acronym_file = 'acronyms_example_sorted.txt'

prs = Presentation()
title_only_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_only_slide_layout)
shapes = slide.shapes

# Count number of items in the acronym file
with open(acronym_file, 'r') as fp:
        lines = len(fp.readlines())

# Parameterize the table dimensions, include header row
rows = lines + 1
cols = 2

left = top = Inches(1.0)
width = Inches(10.0)
height = Inches(0.8)

table = shapes.add_table(rows, cols, left, top, width, height).table

# Set column widths
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(6.0)

# Write column headings
table.cell(0, 0).text = 'Acronym'
table.cell(0, 1).text = 'Description'

# Write acronym cells
row = 1
with open(acronym_file, "rt") as f:
    for line in f:
        x = line.split(',')[0]
        table.cell(row, 0).text = x
        row = row + 1
f.close()

# Write description cells
row = 1
with open(acronym_file, "rt") as f:
    for line in f:
        x = line.split(',')[1]
        table.cell(row, 1).text = x
        row = row + 1
f.close()

prs.save('acronym_table.pptx')
